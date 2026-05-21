"""Yalla Wassel — pitch deck builder.

Produces pitch.pptx — a 10-slide deck with the same content and palette
as pitch.html, but as a native PowerPoint file. Designed for 16:9.

Run:  python build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────
INK = RGBColor(0x11, 0x18, 0x27)
INK_SOFT = RGBColor(0x37, 0x41, 0x51)
INK_MUTED = RGBColor(0x6B, 0x72, 0x80)
INK_FAINT = RGBColor(0x9C, 0xA3, 0xAF)
TRUST = RGBColor(0x7C, 0x3A, 0xED)
TRUST_DARK = RGBColor(0x5B, 0x21, 0xB6)
TRUST_DEEP = RGBColor(0x4C, 0x1D, 0x95)
TRUST_LIGHT = RGBColor(0xA8, 0x55, 0xF7)
TRUST_SOFT = RGBColor(0xED, 0xE9, 0xFE)
TRUST_SOFTER = RGBColor(0xF5, 0xF3, 0xFF)
SURFACE_LINE = RGBColor(0xEF, 0xEC, 0xF6)
SURFACE_BG = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return s


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.75, corner=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # Tune corner radius via adjustment value (0..0.5 of shorter side)
    shape.adjustments[0] = corner
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text,
    *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP, letter_spacing=0, font="Inter",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # Normalize input. `text` can be:
    #   - a str (optionally multi-line via "\n")
    #   - a list of (segment_text, segment_kwargs) tuples → single line of mixed runs
    #   - a list of paragraphs where each paragraph is itself a str or a list of segments
    if isinstance(text, str):
        lines = text.split("\n")
    elif isinstance(text, list):
        if text and isinstance(text[0], tuple):
            lines = [text]  # single line, multiple styled runs
        else:
            lines = text  # list of paragraphs
    else:
        lines = [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, list):
            # list of (segment_text, segment_kwargs) tuples
            for j, seg in enumerate(line):
                seg_text, seg_kwargs = seg
                run = p.add_run()
                run.text = seg_text
                _apply_run_style(run, **{**dict(size=size, bold=bold, color=color, font=font, letter_spacing=letter_spacing), **seg_kwargs})
        else:
            run = p.add_run()
            run.text = line
            _apply_run_style(run, size=size, bold=bold, color=color, font=font, letter_spacing=letter_spacing)
    return tb


def _apply_run_style(run, *, size, bold, color, font, letter_spacing=0):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.color.rgb = color
    if letter_spacing:
        # PPTX uses 'spc' attribute in 100ths of a point
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))


def add_eyebrow(slide, x, y, text, color=TRUST):
    add_text(slide, x, y, Inches(6), Inches(0.32), text.upper(),
             size=10.5, bold=True, color=color, letter_spacing=2.4)


def add_pill(slide, x, y, label):
    """Small rounded pill — used for the 'Pillar 01' tags."""
    w = Inches(0.95)
    h = Inches(0.32)
    add_rect(slide, x, y, w, h, fill=TRUST_SOFT, corner=0.5)
    tb = add_text(slide, x, y, w, h, label.upper(),
                  size=9, bold=True, color=TRUST_DARK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  letter_spacing=1.5)
    return tb


def add_eyebrow_pill(slide, x, y, label):
    """Eyebrow on a soft pill — used at the top of most slides."""
    w = Inches(1.95)
    h = Inches(0.36)
    add_rect(slide, x, y, w, h, fill=TRUST_SOFTER, line=TRUST_SOFT, corner=0.5)
    add_text(slide, x, y, w, h, label.upper(),
             size=10, bold=True, color=TRUST_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             letter_spacing=2.5)


def add_meta_corner(slide):
    """Brand wordmark in the top-left corner."""
    # mark — small purple square
    add_rect(slide, Inches(0.45), Inches(0.42), Inches(0.32), Inches(0.32),
             fill=TRUST, corner=0.32)
    add_text(slide, Inches(0.85), Inches(0.42), Inches(3), Inches(0.32),
             [[("Yalla ", {"color": INK, "bold": True, "size": 11}),
               ("Wassel", {"color": TRUST, "bold": True, "size": 11})]],
             anchor=MSO_ANCHOR.MIDDLE)


def add_progress(slide, idx, total=10):
    """Top-edge progress bar."""
    bar_w = SLIDE_W * (idx / total)
    add_rect(slide, 0, 0, bar_w, Emu(38100), fill=TRUST_LIGHT, corner=0)


def add_page_num(slide, idx, total=10):
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.3),
             f"{idx} / {total}", size=10, bold=True, color=INK_MUTED,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 1)
add_meta_corner(s)

# Soft violet decorative blob (top-right) — emulated with a rounded shape
blob = add_rect(s, Inches(8.5), Inches(-1.2), Inches(6), Inches(6),
                fill=TRUST_SOFTER, corner=0.5)
blob2 = add_rect(s, Inches(-1.5), Inches(5.2), Inches(4), Inches(4),
                 fill=TRUST_SOFTER, corner=0.5)

# Eyebrow
add_eyebrow_pill(s, Inches(0.85), Inches(1.7), "Amman · 2026")

# Headline — three lines, "without" in violet
add_text(s, Inches(0.85), Inches(2.25), Inches(11.5), Inches(3.5),
         "Accountability", size=72, bold=True, color=INK)
add_text(s, Inches(0.85), Inches(3.2), Inches(11.5), Inches(1.2),
         "without", size=72, bold=True, color=TRUST_DARK)
add_text(s, Inches(0.85), Inches(4.15), Inches(11.5), Inches(1.2),
         "surveillance.", size=72, bold=True, color=INK)

# Subhead
add_text(s, Inches(0.85), Inches(5.45), Inches(11), Inches(0.6),
         "A same-day delivery operation that ditched GPS — and stopped losing drivers.",
         size=20, color=INK_SOFT)

# Brand tag
add_text(s, Inches(0.85), Inches(6.2), Inches(11), Inches(0.4),
         "Yalla Wassel — built for Hadeel's team in Wadi Saqra.",
         size=14, bold=True, color=TRUST)

add_page_num(s, 1)

# ────────────────────────────────────────────────────────────────────
# Slide 2 — The Problem
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 2)
add_meta_corner(s)

add_eyebrow_pill(s, Inches(0.85), Inches(1.4), "The problem")

# Left column
add_text(s, Inches(0.85), Inches(1.95), Inches(6.5), Inches(1.4),
         "Eight drivers.\nOne sticky-note operation.",
         size=42, bold=True, color=INK)
add_text(s, Inches(0.85), Inches(4.0), Inches(6.5), Inches(2),
         "Hadeel runs a same-day delivery shop. Orders arrive via WhatsApp, "
         "yelling, and Post-its. She installed GPS trackers to bring order.",
         size=16, color=INK_SOFT)
add_text(s, Inches(0.85), Inches(5.6), Inches(6.5), Inches(0.7),
         "Two drivers quit the same week. One said it out loud.",
         size=16, color=INK_MUTED)

# Right column — quote card
qx, qy, qw, qh = Inches(8.0), Inches(2.2), Inches(4.5), Inches(3.4)
add_rect(s, qx, qy, qw, qh, fill=TRUST_SOFTER, line=TRUST_SOFT, corner=0.12)
add_text(s, qx + Inches(0.4), qy + Inches(0.3), Inches(0.4), Inches(0.8),
         "“", size=64, bold=True, color=TRUST)
add_text(s, qx + Inches(0.4), qy + Inches(0.9), qw - Inches(0.8), Inches(2.0),
         "You don't pay me enough to be watched all day.",
         size=24, bold=True, color=INK)
add_text(s, qx + Inches(0.4), qy + Inches(2.8), qw - Inches(0.8), Inches(0.4),
         "— Mahmoud, before resigning", size=12, color=INK_MUTED)

add_page_num(s, 2)

# ────────────────────────────────────────────────────────────────────
# Slide 3 — The Insight
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 3)
add_meta_corner(s)

add_eyebrow_pill(s, Inches(5.55), Inches(0.95), "The insight")

# Centered headline
add_text(s, Inches(1), Inches(1.45), Inches(11.3), Inches(1.4),
         [[("Surveillance and ", {"color": INK, "size": 40, "bold": True}),
           ("accountability ", {"color": TRUST_DARK, "size": 40, "bold": True}),
           ("are different problems.", {"color": INK, "size": 40, "bold": True})]],
         align=PP_ALIGN.CENTER, size=40, bold=True)

add_text(s, Inches(2), Inches(3.0), Inches(9.3), Inches(0.9),
         "Accountability is the result of clear commitments, not constant observation.",
         size=18, color=INK_SOFT, align=PP_ALIGN.CENTER)

# Two columns comparing
def compare_col(x, y, w, h, title, items, *, accent):
    add_rect(s, x, y, w, h,
             fill=TRUST_SOFTER if accent else SURFACE_BG,
             line=TRUST_SOFT if accent else SURFACE_LINE,
             corner=0.08)
    add_text(s, x + Inches(0.4), y + Inches(0.35), w - Inches(0.8), Inches(0.35),
             title.upper(), size=10.5, bold=True,
             color=TRUST_DARK if accent else INK_FAINT, letter_spacing=2.4)
    for i, item in enumerate(items):
        add_text(s, x + Inches(0.4), y + Inches(0.85) + Inches(i * 0.5),
                 w - Inches(0.8), Inches(0.45),
                 f"— {item}", size=14,
                 color=INK if accent else INK_MUTED)

compare_col(Inches(1.2), Inches(4.0), Inches(5.2), Inches(2.7),
            "GPS tells you",
            ["Coordinates every 30s", "Speed and idle time", "Route deviations"],
            accent=False)
compare_col(Inches(6.9), Inches(4.0), Inches(5.2), Inches(2.7),
            "What actually matters",
            ["Did they pick it up?", "Are they on their way?", "Did the customer get it?"],
            accent=True)

add_page_num(s, 3)

# ────────────────────────────────────────────────────────────────────
# Slide 4 — The Solution
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 4)
add_meta_corner(s)

add_eyebrow_pill(s, Inches(0.85), Inches(1.4), "The solution")
add_text(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(1.1),
         "Yalla Wassel runs on three primitives.",
         size=44, bold=True, color=INK)
add_text(s, Inches(0.85), Inches(3.05), Inches(11.5), Inches(0.5),
         "Not surveillance. Three small ideas, working together.",
         size=18, color=INK_MUTED)

# Three pillar cards
def pillar_card(x, y, n, title, body):
    w, h = Inches(3.85), Inches(3.0)
    add_rect(s, x, y, w, h, fill=WHITE, line=SURFACE_LINE, corner=0.10)
    add_text(s, x + Inches(0.4), y + Inches(0.45), w - Inches(0.8), Inches(0.35),
             n, size=11, bold=True, color=TRUST, letter_spacing=1.5,
             font="Consolas")
    add_text(s, x + Inches(0.4), y + Inches(0.95), w - Inches(0.8), Inches(0.8),
             title, size=22, bold=True, color=INK)
    add_text(s, x + Inches(0.4), y + Inches(1.7), w - Inches(0.8), Inches(1.2),
             body, size=14, color=INK_MUTED)

pillar_card(Inches(0.85), Inches(3.85), "01", "Trust Ledger",
            "Four checkpoint taps per delivery. Append-only, public commitments. Zero coordinates stored.")
pillar_card(Inches(4.85), Inches(3.85), "02", "Mutual Trust Score",
            "Drivers rate the system every week. The dispatcher sees their score next to her own.")
pillar_card(Inches(8.85), Inches(3.85), "03", "Explainable Dispatch",
            "Every assignment ships with one sentence. Override the engine? You have to say why.")

add_page_num(s, 4)

# ────────────────────────────────────────────────────────────────────
# Slide 5 — Pillar 1: Trust Ledger
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 5)
add_meta_corner(s)

add_pill(s, Inches(0.85), Inches(1.4), "Pillar 01")

add_text(s, Inches(0.85), Inches(1.95), Inches(11), Inches(1),
         "The Trust Ledger.", size=46, bold=True, color=INK)
add_text(s, Inches(0.85), Inches(3.05), Inches(11), Inches(1.1),
         "Drivers tap four buttons per delivery. Each tap is a public commitment, "
         "not a tracked coordinate. Append-only, auditable, dignified.",
         size=17, color=INK_SOFT)

# Checkpoint pills
def checkpoint(x, y, label, time, on):
    w, h = Inches(2.4), Inches(1.4)
    add_rect(s, x, y, w, h,
             fill=TRUST_SOFTER if on else WHITE,
             line=TRUST_LIGHT if on else SURFACE_LINE,
             corner=0.16)
    # dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             x + Inches(0.35), y + Inches(0.35),
                             Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = TRUST if on else INK_FAINT
    dot.line.fill.background()
    dot.shadow.inherit = False

    add_text(s, x + Inches(0.35), y + Inches(0.7), w - Inches(0.7), Inches(0.4),
             label, size=14, bold=True, color=INK if on else INK_MUTED)
    add_text(s, x + Inches(0.35), y + Inches(1.0), w - Inches(0.7), Inches(0.3),
             time, size=11, color=INK_MUTED if on else INK_FAINT)

def arrow(x, y):
    add_text(s, x, y, Inches(0.5), Inches(1.4),
             "→", size=22, color=TRUST_LIGHT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)

# Layout: 4 checkpoints + 3 arrows
base_y = Inches(4.6)
checkpoint(Inches(0.85), base_y, "Assigned", "08:14", True)
arrow(Inches(3.3), base_y)
checkpoint(Inches(3.85), base_y, "Picked up", "08:28", True)
arrow(Inches(6.3), base_y)
checkpoint(Inches(6.85), base_y, "Arrived nearby", "08:41", True)
arrow(Inches(9.3), base_y)
checkpoint(Inches(9.85), base_y, "Delivered", "—", False)

# Subtle note
add_text(s, Inches(0.85), Inches(6.4), Inches(11.5), Inches(0.6),
         "Optional delay reasons live alongside — traffic, store delay, customer unavailable, vehicle. One tap. Customer sees it.",
         size=12, color=INK_MUTED)

add_page_num(s, 5)

# ────────────────────────────────────────────────────────────────────
# Slide 6 — Pillar 2: Mutual Trust
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 6)
add_meta_corner(s)

add_pill(s, Inches(0.85), Inches(1.4), "Pillar 02")

add_text(s, Inches(0.85), Inches(1.95), Inches(11), Inches(1.1),
         [[("Trust runs ", {"color": INK, "size": 46, "bold": True}),
           ("both ways.", {"color": TRUST_DARK, "size": 46, "bold": True})]],
         size=46, bold=True)
add_text(s, Inches(0.85), Inches(3.05), Inches(11.5), Inches(0.8),
         "Drivers rate the system weekly on three axes: fairness, fit, pressure. "
         "The dispatcher sees their team's score next to her own.",
         size=17, color=INK_SOFT)

# Three score cards
def score_card(x, y, num, label, sub, *, hero=False):
    w, h = Inches(3.85), Inches(2.7)
    add_rect(s, x, y, w, h,
             fill=TRUST_SOFTER if hero else WHITE,
             line=TRUST_SOFT if hero else SURFACE_LINE,
             corner=0.10)
    add_text(s, x, y + Inches(0.4), w, Inches(1.4),
             str(num), size=72, bold=True,
             color=TRUST_DARK if hero else INK,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.85), w, Inches(0.4),
             label, size=15, bold=True,
             color=TRUST_DARK if hero else INK,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.25), w, Inches(0.3),
             sub, size=11, color=INK_MUTED, align=PP_ALIGN.CENTER)

base_y = Inches(4.1)
score_card(Inches(0.85), base_y, 88, "Driver score", "on-time, completion, streaks")
score_card(Inches(4.85), base_y, 87, "Mutual Trust Index", "geometric mean — both sides count", hero=True)
score_card(Inches(8.85), base_y, 85, "System score", "fairness, fit, pressure")

add_text(s, Inches(0.85), Inches(6.95), Inches(11.5), Inches(0.4),
         "Geometric mean — you can't trade one side for the other.",
         size=12, color=INK_MUTED)

add_page_num(s, 6)

# ────────────────────────────────────────────────────────────────────
# Slide 7 — Pillar 3: Explainable Dispatch
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 7)
add_meta_corner(s)

add_pill(s, Inches(0.85), Inches(1.4), "Pillar 03")

add_text(s, Inches(0.85), Inches(1.95), Inches(11), Inches(1.1),
         "Every assignment, one sentence.", size=44, bold=True, color=INK)

# Suggestion card
cx, cy, cw, ch = Inches(0.85), Inches(3.2), Inches(7.5), Inches(3.5)
add_rect(s, cx, cy, cw, ch, fill=TRUST_SOFTER, line=TRUST_SOFT, corner=0.10)
add_text(s, cx + Inches(0.4), cy + Inches(0.35), cw - Inches(0.8), Inches(0.3),
         "ENGINE SUGGESTS · #1005 · URGENT",
         size=10.5, bold=True, color=TRUST, letter_spacing=2.2)
add_text(s, cx + Inches(0.4), cy + Inches(0.8), cw - Inches(0.8), Inches(0.7),
         "Hamza — Central", size=26, bold=True, color=INK)

reasons = [
    "Jabal Amman is in his coverage",
    "0 active orders right now",
    "96% on-time rate on urgent",
    "Last 4 deliveries to this street were his",
]
for i, r in enumerate(reasons):
    y = cy + Inches(1.7) + Inches(i * 0.42)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             cx + Inches(0.4), y + Inches(0.13),
                             Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = TRUST
    dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s, cx + Inches(0.65), y, cw - Inches(1.1), Inches(0.4),
             r, size=14, color=INK_SOFT)

# Right-side note
add_text(s, Inches(8.7), Inches(3.4), Inches(3.9), Inches(0.6),
         "Take the pick.", size=22, bold=True, color=INK)
add_text(s, Inches(8.7), Inches(4.0), Inches(3.9), Inches(0.6),
         "One click.", size=22, bold=True, color=INK)

add_text(s, Inches(8.7), Inches(5.0), Inches(3.9), Inches(1.6),
         "Override? You have to say why. That reason becomes a signal "
         "the engine learns from.",
         size=14, color=INK_MUTED)

add_page_num(s, 7)

# ────────────────────────────────────────────────────────────────────
# Slide 8 — Three Roles
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 8)
add_meta_corner(s)

add_eyebrow_pill(s, Inches(0.85), Inches(1.4), "Three roles")
add_text(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(1.1),
         "Built around the way work actually happens.",
         size=40, bold=True, color=INK)

def role_card(x, y, tag, name, body, *, hero=False):
    w, h = Inches(3.85), Inches(3.2)
    add_rect(s, x, y, w, h,
             fill=TRUST_SOFTER if hero else WHITE,
             line=TRUST_LIGHT if hero else SURFACE_LINE,
             corner=0.10)
    add_text(s, x + Inches(0.4), y + Inches(0.4), w - Inches(0.8), Inches(0.3),
             tag.upper(), size=10.5, bold=True, color=TRUST, letter_spacing=2)
    add_text(s, x + Inches(0.4), y + Inches(0.95), w - Inches(0.8), Inches(0.6),
             name, size=22, bold=True, color=INK)
    add_text(s, x + Inches(0.4), y + Inches(1.7), w - Inches(0.8), Inches(1.4),
             body, size=14, color=INK_MUTED)

role_card(Inches(0.85), Inches(3.35), "Dispatcher", "Command center",
          "Live timeline, explainable dispatch, mutual-trust dashboard. No map.")
role_card(Inches(4.85), Inches(3.35), "Driver", "Silent Mode",
          "One screen, one order, one button at a time. Zero notifications.",
          hero=True)
role_card(Inches(8.85), Inches(3.35), "Customer", "Live tracking",
          "Timeline, not a map. ETA, status, confirmation code. Never sees driver location.")

add_text(s, Inches(0.85), Inches(6.8), Inches(11.5), Inches(0.4),
         [[("The dispatcher screen has ", {"color": INK_MUTED, "size": 13}),
           ("no map", {"color": TRUST_DARK, "size": 13, "bold": True}),
           (". That deliberate absence is the entire product thesis.", {"color": INK_MUTED, "size": 13})]])

add_page_num(s, 8)

# ────────────────────────────────────────────────────────────────────
# Slide 9 — Why We Win
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 9)
add_meta_corner(s)

add_eyebrow_pill(s, Inches(0.85), Inches(1.4), "Why we win")
add_text(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(1),
         "The moat is a team that doesn't quit.",
         size=40, bold=True, color=INK)

def compare_card(x, y, label, body):
    w, h = Inches(3.85), Inches(2.4)
    add_rect(s, x, y, w, h, fill=WHITE, line=SURFACE_LINE, corner=0.10)
    add_text(s, x + Inches(0.4), y + Inches(0.4), w - Inches(0.8), Inches(0.3),
             label.upper(), size=10.5, bold=True, color=TRUST, letter_spacing=1.6,
             font="Consolas")
    add_text(s, x + Inches(0.4), y + Inches(0.95), w - Inches(0.8), Inches(1.3),
             body, size=14, color=INK)

base_y = Inches(3.25)
compare_card(Inches(0.85), base_y, "vs Onfleet / Bringg",
             "Built for fleets of interchangeable contractors. We build for retention.")
compare_card(Inches(4.85), base_y, "vs WhatsApp + sheets",
             "10× more legible without becoming 10× more invasive.")
compare_card(Inches(8.85), base_y, "vs Full GPS surveillance",
             "Surveillance has a hidden cost: the best people leave.")

# Gradient CTA card (north-star)
nx, ny, nw, nh = Inches(0.85), Inches(5.9), Inches(11.65), Inches(1.0)
add_rect(s, nx, ny, nw, nh, fill=TRUST_DARK, corner=0.18)
add_text(s, nx + Inches(0.5), ny, nw - Inches(1.0), nh,
         [[("North-star metric:  ", {"color": WHITE, "size": 16}),
           ("Mutual Trust Index ≥ 85", {"color": WHITE, "size": 18, "bold": True}),
           ("  within month one. Track the team you can't afford to lose.",
            {"color": WHITE, "size": 16})]],
         anchor=MSO_ANCHOR.MIDDLE)

add_page_num(s, 9)

# ────────────────────────────────────────────────────────────────────
# Slide 10 — Closing
# ────────────────────────────────────────────────────────────────────

s = add_slide()
add_progress(s, 10)
add_meta_corner(s)

add_eyebrow_pill(s, Inches(5.5), Inches(1.45), "What's next")

# Big closing line, centered, two halves with gradient word
add_text(s, Inches(1), Inches(2.4), Inches(11.3), Inches(1.0),
         [[("The dispatcher screen has ", {"color": INK, "size": 40, "bold": True}),
           ("no map.", {"color": TRUST_DARK, "size": 40, "bold": True})]],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.35), Inches(11.3), Inches(0.9),
         "That deliberate absence is the entire product.",
         size=40, bold=True, color=INK, align=PP_ALIGN.CENTER)

# Roadmap
add_text(s, Inches(1.5), Inches(4.85), Inches(10.3), Inches(1.2),
         "v0.1 ships today. v0.2 onboards 5 more Amman operations. "
         "v1.0 opens the checkpoint primitive as an API — trust infrastructure "
         "for any work where people feel surveilled.",
         size=15, color=INK_MUTED, align=PP_ALIGN.CENTER)

# CTA — gradient button
cx, cy = Inches(4.8), Inches(6.0)
cw, ch = Inches(3.7), Inches(0.7)
add_rect(s, cx, cy, cw, ch, fill=TRUST_DARK, corner=0.42)
add_text(s, cx, cy, cw, ch,
         "github.com/anasmutlaq123-ship-it/yalla-wassel",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.85), Inches(7.0), Inches(11.5), Inches(0.4),
         "Yalla Wassel · Amman · since 2024",
         size=11, color=INK_MUTED, align=PP_ALIGN.CENTER)

add_page_num(s, 10)

# ────────────────────────────────────────────────────────────────────

out = "pitch.pptx"
prs.save(out)
print(f"Wrote {out} — {len(prs.slides)} slides")
